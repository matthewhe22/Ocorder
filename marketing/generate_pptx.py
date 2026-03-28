"""
Generate TOCS OC Order Portal Marketing Plan PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Color palette ──────────────────────────────────────────────────────────────
FOREST_GREEN  = RGBColor(0x1B, 0x3F, 0x6B)   # deep navy blue
SAGE_GREEN    = RGBColor(0x5B, 0x9B, 0xD5)   # sky blue
CREAM         = RGBColor(0xEB, 0xF5, 0xFB)   # very light blue
CHARCOAL      = RGBColor(0x1A, 0x2A, 0x3A)   # dark blue-charcoal
WHITE         = RGBColor(0xE0, 0xF0, 0xFF)   # light blue (slide backgrounds)
LIGHT_SAGE    = RGBColor(0xC5, 0xDC, 0xF0)   # pale blue
MID_GREEN     = RGBColor(0x2E, 0x6D, 0xA4)   # medium blue

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helper utilities ───────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # MSO_SHAPE.RECTANGLE = 1
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                font_name="Calibri", italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, font_name="Calibri", space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = color
    return p


def slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def section_title(slide, title_text, fg=WHITE, bg=FOREST_GREEN,
                  top=Inches(0), height=Inches(1.15)):
    rect = add_rect(slide, Inches(0), top, SLIDE_W, height, fill_color=bg)
    add_textbox(slide, Inches(0.4), top + Inches(0.15), Inches(12.5), height - Inches(0.1),
                title_text, font_size=32, bold=True, color=fg, align=PP_ALIGN.LEFT)
    return rect


def green_card(slide, left, top, width, height, title, body,
               title_color=WHITE, body_color=CHARCOAL,
               card_bg=FOREST_GREEN, body_bg=CREAM):
    # title bar
    add_rect(slide, left, top, width, Inches(0.55), fill_color=card_bg)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.05),
                width - Inches(0.2), Inches(0.5),
                title, font_size=13, bold=True, color=title_color, align=PP_ALIGN.LEFT)
    # body
    body_top = top + Inches(0.55)
    body_h = height - Inches(0.55)
    add_rect(slide, left, body_top, width, body_h, fill_color=body_bg)
    add_textbox(slide, left + Inches(0.12), body_top + Inches(0.1),
                width - Inches(0.24), body_h - Inches(0.15),
                body, font_size=11, color=body_color, align=PP_ALIGN.LEFT, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
slide_bg(s1, FOREST_GREEN)

# Decorative horizontal band
add_rect(s1, Inches(0), Inches(2.8), SLIDE_W, Inches(0.08), fill_color=SAGE_GREEN)
add_rect(s1, Inches(0), Inches(5.2), SLIDE_W, Inches(0.05), fill_color=SAGE_GREEN)

# Logo-style block top-left
add_rect(s1, Inches(0.5), Inches(0.45), Inches(1.1), Inches(0.55), fill_color=SAGE_GREEN)
add_textbox(s1, Inches(0.52), Inches(0.47), Inches(1.06), Inches(0.52),
            "TOCS", font_size=20, bold=True, color=FOREST_GREEN, align=PP_ALIGN.CENTER)

# Main title
add_textbox(s1, Inches(0.5), Inches(1.1), Inches(12.0), Inches(1.6),
            "TOCS OC Order Portal", font_size=52, bold=True, color=WHITE,
            align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(s1, Inches(0.5), Inches(2.9), Inches(12.0), Inches(0.8),
            "OC Certificates. Online. On Time.", font_size=26, bold=False,
            color=CREAM, align=PP_ALIGN.LEFT, italic=True)

# Description line
add_textbox(s1, Inches(0.5), Inches(3.8), Inches(12.0), Inches(0.7),
            "Australia's self-service OC certificate ordering platform for conveyancers,",
            font_size=16, color=LIGHT_SAGE, align=PP_ALIGN.LEFT)
add_textbox(s1, Inches(0.5), Inches(4.3), Inches(12.0), Inches(0.5),
            "solicitors, real estate agents and property managers.",
            font_size=16, color=LIGHT_SAGE, align=PP_ALIGN.LEFT)

# Footer bar
add_rect(s1, Inches(0), Inches(6.9), SLIDE_W, Inches(0.6), fill_color=MID_GREEN)
add_textbox(s1, Inches(0.5), Inches(6.92), Inches(12.0), Inches(0.5),
            "Marketing Plan 2026  |  TOCS — Top Owners Corporation Solution",
            font_size=12, color=CREAM, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
slide_bg(s2, WHITE)
section_title(s2, "The Problem We Solve")

pain_points = [
    ("☎  Phone Calls & Long Waits",
     "Ordering OC certificates requires phone calls, hold times, and back-and-forth emails with strata managers. Time-critical settlements are put at risk."),
    ("📄  Manual Forms, Fax & Email",
     "Paper-based and email workflows are slow, error-prone and hard to track. Staff spend hours chasing confirmations and correcting mistakes."),
    ("🔍  No Visibility on Order Status",
     "Once an order is placed, customers are in the dark. There's no status tracking, no estimated delivery, and no automated notifications."),
]

card_w = Inches(3.9)
card_h = Inches(4.6)
tops   = Inches(1.4)
gap    = Inches(0.37)

for i, (title, body) in enumerate(pain_points):
    left = Inches(0.42) + i * (card_w + gap)
    green_card(s2, left, tops, card_w, card_h,
               title, body,
               card_bg=FOREST_GREEN, body_bg=LIGHT_SAGE,
               title_color=WHITE, body_color=CHARCOAL)

# Bottom accent
add_rect(s2, Inches(0), Inches(6.85), SLIDE_W, Inches(0.15), fill_color=SAGE_GREEN)
add_textbox(s2, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
            "TOCS OC Order Portal  |  Marketing Plan 2026",
            font_size=10, color=CHARCOAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
slide_bg(s3, WHITE)
section_title(s3, "Introducing the TOCS OC Order Portal")

add_textbox(s3, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.5),
            "Australia's first fully self-service OC certificate ordering platform",
            font_size=16, italic=True, color=MID_GREEN, align=PP_ALIGN.LEFT)

features = [
    ("🕐  Order Online 24/7",
     "No phone calls needed. Place orders any time from any device. Confirmation is instant."),
    ("📡  Real-Time Tracking",
     "Know your order status at every step. Automated notifications keep you and your client informed."),
    ("💳  Multiple Payment Options",
     "Card, PayID, Bank Transfer, or Invoice. Flexible billing to suit every firm."),
    ("📧  Secure Document Delivery",
     "Certificates are delivered directly to your inbox as secure PDFs. No chasing, no delays."),
]

card_w2 = Inches(2.95)
card_h2 = Inches(4.0)
top2    = Inches(1.85)
gap2    = Inches(0.28)

for i, (title, body) in enumerate(features):
    left = Inches(0.42) + i * (card_w2 + gap2)
    green_card(s3, left, top2, card_w2, card_h2,
               title, body,
               card_bg=FOREST_GREEN, body_bg=CREAM,
               title_color=WHITE, body_color=CHARCOAL)

add_rect(s3, Inches(0), Inches(6.85), SLIDE_W, Inches(0.15), fill_color=SAGE_GREEN)
add_textbox(s3, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
            "TOCS OC Order Portal  |  Marketing Plan 2026",
            font_size=10, color=CHARCOAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Target Market (2×2 grid)
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
slide_bg(s4, WHITE)
section_title(s4, "Who Uses the Platform")

segments = [
    ("⚖  Conveyancers & Solicitors",
     "The primary users. They order OC certificates for every property settlement. Speed and accuracy are critical — delays cost money and client trust."),
    ("🏠  Real Estate Agents",
     "Agents need OC documents for listings and sales. A self-service portal removes their dependency on strata managers and speeds up transactions."),
    ("👤  Property Owners",
     "Owners selling or refinancing need OC certificates. A simple online portal lets them order directly without intermediaries."),
    ("🏢  Building Managers",
     "Strata and building managers use the portal for bulk ordering, compliance certificates, and managing requests across multiple properties."),
]

c_w = Inches(6.0)
c_h = Inches(2.4)
positions = [
    (Inches(0.42), Inches(1.35)),
    (Inches(6.85), Inches(1.35)),
    (Inches(0.42), Inches(3.95)),
    (Inches(6.85), Inches(3.95)),
]

for (left, top), (title, body) in zip(positions, segments):
    green_card(s4, left, top, c_w, c_h,
               title, body,
               card_bg=FOREST_GREEN, body_bg=LIGHT_SAGE,
               title_color=WHITE, body_color=CHARCOAL)

add_rect(s4, Inches(0), Inches(6.85), SLIDE_W, Inches(0.15), fill_color=SAGE_GREEN)
add_textbox(s4, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
            "TOCS OC Order Portal  |  Marketing Plan 2026",
            font_size=10, color=CHARCOAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Product & Pricing (table)
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
slide_bg(s5, WHITE)
section_title(s5, "What You Can Order")

rows = [
    ("Product / Service",         "Standard",   "Urgent / Express"),
    ("OC Certificate",            "$220",        "$385"),
    ("Building Register Search",  "POA",         "—"),
    ("Insurance Certificate of Currency", "POA", "—"),
    ("Meeting Minutes & Financial Statements", "POA", "—"),
    ("Keys / Fobs / Remotes",     "Per item",    "—"),
    ("Shipping",                  "$15",         "$25 express"),
]

tbl_left   = Inches(1.0)
tbl_top    = Inches(1.3)
tbl_width  = Inches(11.3)
tbl_height = Inches(5.0)

table = s5.shapes.add_table(len(rows), 3, tbl_left, tbl_top, tbl_width, tbl_height).table

col_widths = [Inches(6.2), Inches(2.5), Inches(2.6)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

for r_idx, row_data in enumerate(rows):
    for c_idx, cell_text in enumerate(row_data):
        cell = table.cell(r_idx, c_idx)
        cell.text = cell_text
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
        run.text = cell_text
        run.font.name = "Calibri"
        run.font.size = Pt(13 if r_idx == 0 else 12)
        run.font.bold = (r_idx == 0)
        if r_idx == 0:
            run.font.color.rgb = WHITE
        elif r_idx % 2 == 0:
            run.font.color.rgb = CHARCOAL
        else:
            run.font.color.rgb = CHARCOAL

        # Cell fill
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        # Remove existing solidFill if present
        for old in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(old)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        if r_idx == 0:
            srgb.set('val', '2D5016')
        elif r_idx % 2 == 0:
            srgb.set('val', 'F5F0E8')
        else:
            srgb.set('val', 'D6E8C8')

add_rect(s5, Inches(0), Inches(6.85), SLIDE_W, Inches(0.15), fill_color=SAGE_GREEN)
add_textbox(s5, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
            "TOCS OC Order Portal  |  Marketing Plan 2026",
            font_size=10, color=CHARCOAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Market Opportunity
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
slide_bg(s6, FOREST_GREEN)
section_title(s6, "The Market Opportunity", fg=WHITE, bg=MID_GREEN)

stats = [
    ("320,000+",   "strata schemes in Australia"),
    ("100,000s",   "property settlements per year require OC documents"),
    ("~90%",       "of OC orders still placed manually — phone & email"),
    ("↑ Growing",  "Fast-growing apartment market driving demand for digital solutions"),
]

stat_w = Inches(5.8)
stat_h = Inches(2.2)
positions6 = [
    (Inches(0.42), Inches(1.45)),
    (Inches(7.0),  Inches(1.45)),
    (Inches(0.42), Inches(3.85)),
    (Inches(7.0),  Inches(3.85)),
]

for (left, top), (number, label) in zip(positions6, stats):
    add_rect(s6, left, top, stat_w, stat_h, fill_color=MID_GREEN)
    add_textbox(s6, left + Inches(0.2), top + Inches(0.15),
                stat_w - Inches(0.3), Inches(0.75),
                number, font_size=34, bold=True, color=CREAM, align=PP_ALIGN.LEFT)
    add_textbox(s6, left + Inches(0.2), top + Inches(0.9),
                stat_w - Inches(0.3), stat_h - Inches(0.95),
                label, font_size=14, color=LIGHT_SAGE, align=PP_ALIGN.LEFT, wrap=True)

add_rect(s6, Inches(0), Inches(6.85), SLIDE_W, Inches(0.15), fill_color=SAGE_GREEN)
add_textbox(s6, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
            "TOCS OC Order Portal  |  Marketing Plan 2026",
            font_size=10, color=CREAM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Launch Campaign
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
slide_bg(s7, WHITE)
section_title(s7, "Launch Campaign: Order Smarter")

add_textbox(s7, Inches(0.4), Inches(1.22), Inches(12.5), Inches(0.45),
            "Campaign Theme: \"OC documents shouldn't slow down your settlement\"",
            font_size=15, italic=True, color=MID_GREEN, align=PP_ALIGN.LEFT)

phases = [
    ("Phase 1",  "Weeks 1–4",   "Soft Launch",
     "Onboard existing TOCS customers. Gather feedback, refine UX, build case studies."),
    ("Phase 2",  "Weeks 5–8",   "Industry Seeding",
     "Engage AIC, SCA, REIA. Post on LinkedIn. Reach out to conveyancing firms directly."),
    ("Phase 3",  "Weeks 9–16",  "Paid Acquisition",
     "Launch Google Ads (high-intent keywords). Retargeting campaigns. Drive trial sign-ups."),
    ("Phase 4",  "Month 5+",    "Scale & Optimise",
     "Analyse performance data. Scale top channels. Launch volume pricing and firm partnerships."),
]

p_w = Inches(3.0)
p_h = Inches(3.8)
top7 = Inches(1.85)
gap7 = Inches(0.22)

colors7 = [FOREST_GREEN, MID_GREEN, SAGE_GREEN, LIGHT_SAGE]
txt_colors7 = [WHITE, WHITE, CHARCOAL, CHARCOAL]

for i, (phase, timing, headline, detail) in enumerate(phases):
    left = Inches(0.42) + i * (p_w + gap7)
    add_rect(s7, left, top7, p_w, p_h, fill_color=colors7[i])
    add_textbox(s7, left + Inches(0.1), top7 + Inches(0.1),
                p_w - Inches(0.15), Inches(0.45),
                phase, font_size=13, bold=True, color=txt_colors7[i])
    add_textbox(s7, left + Inches(0.1), top7 + Inches(0.55),
                p_w - Inches(0.15), Inches(0.38),
                timing, font_size=12, italic=True, color=txt_colors7[i])
    add_rect(s7, left + Inches(0.1), top7 + Inches(0.93),
             p_w - Inches(0.2), Inches(0.03), fill_color=CREAM if i < 2 else FOREST_GREEN)
    add_textbox(s7, left + Inches(0.1), top7 + Inches(1.02),
                p_w - Inches(0.15), Inches(0.5),
                headline, font_size=13, bold=True, color=txt_colors7[i])
    add_textbox(s7, left + Inches(0.1), top7 + Inches(1.6),
                p_w - Inches(0.15), p_h - Inches(1.7),
                detail, font_size=11, color=txt_colors7[i], wrap=True)

add_rect(s7, Inches(0), Inches(6.85), SLIDE_W, Inches(0.15), fill_color=SAGE_GREEN)
add_textbox(s7, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
            "TOCS OC Order Portal  |  Marketing Plan 2026",
            font_size=10, color=CHARCOAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Marketing Channels
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
slide_bg(s8, WHITE)
section_title(s8, "Marketing Channels")

channels = [
    ("🔍  SEO",
     "Target high-intent keywords:\n\"OC certificate order online\"\n\"owners corporation certificate\"\n\"strata certificate VIC\"\n\nBlog content for conveyancers & property lawyers."),
    ("📢  Google Ads",
     "Search campaigns targeting settlement-related queries.\n\nRetargeting campaigns for site visitors who didn't convert.\n\nFocus on metro markets: Melbourne, Sydney, Brisbane."),
    ("💼  LinkedIn",
     "Target conveyancers, property lawyers, strata managers, and real estate principals.\n\nSponsored content + InMail campaigns.\n\nThought leadership articles."),
    ("🤝  Industry Partnerships",
     "Sponsorships & co-marketing with:\n- AIC (Australian Institute of Conveyancers)\n- SCA (Strata Community Association)\n- REIA (Real Estate Institute of Australia)"),
]

ch_w = Inches(2.95)
ch_h = Inches(4.3)
top8 = Inches(1.35)
gap8 = Inches(0.28)

for i, (title, body) in enumerate(channels):
    left = Inches(0.42) + i * (ch_w + gap8)
    green_card(s8, left, top8, ch_w, ch_h,
               title, body,
               card_bg=FOREST_GREEN, body_bg=CREAM,
               title_color=WHITE, body_color=CHARCOAL)

add_rect(s8, Inches(0), Inches(6.85), SLIDE_W, Inches(0.15), fill_color=SAGE_GREEN)
add_textbox(s8, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
            "TOCS OC Order Portal  |  Marketing Plan 2026",
            font_size=10, color=CHARCOAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Budget & KPIs
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
slide_bg(s9, WHITE)
section_title(s9, "Budget & Key Metrics")

# --- LEFT: Budget table ---
add_rect(s9, Inches(0.42), Inches(1.35), Inches(6.0), Inches(0.5), fill_color=FOREST_GREEN)
add_textbox(s9, Inches(0.52), Inches(1.38), Inches(5.8), Inches(0.45),
            "Annual Marketing Budget", font_size=14, bold=True, color=WHITE)

budget_items = [
    ("Google Ads",       "$24,000 / yr"),
    ("LinkedIn Ads",     "$18,000 / yr"),
    ("SEO / Content",    "$12,000 / yr"),
    ("Sponsorships",     "$6,000 / yr"),
    ("Creative / Design","$8,000 / yr"),
    ("Other / Buffer",   "$5,400 / yr"),
    ("TOTAL",            "~$73,400 / yr"),
]

row_h = Inches(0.55)
for i, (item, cost) in enumerate(budget_items):
    top_r = Inches(1.85) + i * row_h
    bg = LIGHT_SAGE if i % 2 == 0 else CREAM
    if item == "TOTAL":
        bg = MID_GREEN
    add_rect(s9, Inches(0.42), top_r, Inches(6.0), row_h, fill_color=bg)
    txt_color = WHITE if item == "TOTAL" else CHARCOAL
    bold_row  = (item == "TOTAL")
    add_textbox(s9, Inches(0.55), top_r + Inches(0.1),
                Inches(3.5), row_h - Inches(0.1),
                item, font_size=12, bold=bold_row, color=txt_color)
    add_textbox(s9, Inches(4.1), top_r + Inches(0.1),
                Inches(2.2), row_h - Inches(0.1),
                cost, font_size=12, bold=bold_row, color=txt_color,
                align=PP_ALIGN.RIGHT)

# --- RIGHT: KPIs ---
add_rect(s9, Inches(7.0), Inches(1.35), Inches(5.9), Inches(0.5), fill_color=FOREST_GREEN)
add_textbox(s9, Inches(7.1), Inches(1.38), Inches(5.7), Inches(0.45),
            "KPI Targets", font_size=14, bold=True, color=WHITE)

kpis = [
    ("10,000",   "unique visitors/month   (by Month 6)"),
    ("250",      "active conveyancer firms   (12 months)"),
    ("60%",      "of orders placed online   (12 months)"),
    ("4.5 / 5+", "Net Promoter Score target"),
    ("$320",     "target customer acquisition cost"),
    ("<48 hrs",  "average order fulfilment time"),
]

for i, (metric, label) in enumerate(kpis):
    top_r = Inches(1.85) + i * row_h
    bg = LIGHT_SAGE if i % 2 == 0 else CREAM
    add_rect(s9, Inches(7.0), top_r, Inches(5.9), row_h, fill_color=bg)
    add_textbox(s9, Inches(7.1), top_r + Inches(0.08),
                Inches(1.5), row_h - Inches(0.08),
                metric, font_size=14, bold=True, color=FOREST_GREEN)
    add_textbox(s9, Inches(8.65), top_r + Inches(0.1),
                Inches(4.0), row_h - Inches(0.1),
                label, font_size=11, color=CHARCOAL)

add_rect(s9, Inches(0), Inches(6.85), SLIDE_W, Inches(0.15), fill_color=SAGE_GREEN)
add_textbox(s9, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
            "TOCS OC Order Portal  |  Marketing Plan 2026",
            font_size=10, color=CHARCOAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 12-Month Roadmap
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
slide_bg(s10, WHITE)
section_title(s10, "12-Month Roadmap")

roadmap = [
    ("Month 1",  "Website live, launch collateral ready, team briefed"),
    ("Month 2",  "Soft launch — existing TOCS customers onboarded"),
    ("Month 3",  "Industry outreach begins, LinkedIn campaign live"),
    ("Month 4",  "Google Ads campaigns launched, first blog articles published"),
    ("Month 5",  "Email drip campaigns, first conveyancer webinar"),
    ("Month 6",  "KPI review checkpoint, optimise ad spend, PR outreach"),
    ("Month 9",  "Volume pricing introduced, firm partnership programme"),
    ("Month 12", "Annual review, Year 2 strategy & budget sign-off"),
]

tbl2_left   = Inches(0.5)
tbl2_top    = Inches(1.3)
tbl2_width  = Inches(12.3)
tbl2_height = Inches(5.35)

table2 = s10.shapes.add_table(len(roadmap) + 1, 2, tbl2_left, tbl2_top, tbl2_width, tbl2_height).table
table2.columns[0].width = Inches(2.0)
table2.columns[1].width = Inches(10.3)

# Header row
for c_idx, hdr in enumerate(["Milestone", "Activity"]):
    cell = table2.cell(0, c_idx)
    cell.text = hdr
    tf = cell.text_frame
    run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
    run.text = hdr
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgb.set('val', '2D5016')

for r_idx, (month, activity) in enumerate(roadmap):
    row_idx = r_idx + 1
    for c_idx, val in enumerate([month, activity]):
        cell = table2.cell(row_idx, c_idx)
        cell.text = val
        tf = cell.text_frame
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
        run.text = val
        run.font.size = Pt(12)
        run.font.name = "Calibri"
        run.font.bold = (c_idx == 0)
        run.font.color.rgb = FOREST_GREEN if c_idx == 0 else CHARCOAL
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER if c_idx == 0 else PP_ALIGN.LEFT
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(old)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', 'F5F0E8' if r_idx % 2 == 0 else 'D6E8C8')

add_rect(s10, Inches(0), Inches(6.85), SLIDE_W, Inches(0.15), fill_color=SAGE_GREEN)
add_textbox(s10, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
            "TOCS OC Order Portal  |  Marketing Plan 2026",
            font_size=10, color=CHARCOAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Call to Action
# ══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
slide_bg(s11, FOREST_GREEN)

# Decorative bands
add_rect(s11, Inches(0), Inches(2.6), SLIDE_W, Inches(0.08), fill_color=SAGE_GREEN)
add_rect(s11, Inches(0), Inches(5.5), SLIDE_W, Inches(0.08), fill_color=SAGE_GREEN)

# Logo block
add_rect(s11, Inches(0.5), Inches(0.45), Inches(1.1), Inches(0.55), fill_color=SAGE_GREEN)
add_textbox(s11, Inches(0.52), Inches(0.47), Inches(1.06), Inches(0.52),
            "TOCS", font_size=20, bold=True, color=FOREST_GREEN, align=PP_ALIGN.CENTER)

add_textbox(s11, Inches(0.5), Inches(1.0), Inches(12.0), Inches(1.5),
            "Get Started Today", font_size=46, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_textbox(s11, Inches(0.5), Inches(2.7), Inches(12.0), Inches(0.7),
            "Register your firm and get your first 3 orders free",
            font_size=22, bold=True, color=CREAM, align=PP_ALIGN.LEFT)

add_textbox(s11, Inches(0.5), Inches(3.5), Inches(12.0), Inches(0.5),
            "Experience Australia's fastest, simplest way to order OC certificates.",
            font_size=15, italic=True, color=LIGHT_SAGE, align=PP_ALIGN.LEFT)

# Contact details box
add_rect(s11, Inches(0.5), Inches(4.15), Inches(7.5), Inches(1.25), fill_color=MID_GREEN)
add_textbox(s11, Inches(0.65), Inches(4.2), Inches(7.2), Inches(0.4),
            "Contact & Sign Up", font_size=13, bold=True, color=CREAM)
add_textbox(s11, Inches(0.65), Inches(4.6), Inches(7.2), Inches(0.35),
            "🌐  portal.tocs.com.au", font_size=14, bold=True, color=WHITE)
add_textbox(s11, Inches(0.65), Inches(4.95), Inches(7.2), Inches(0.35),
            "📧  hello@tocs.com.au  |  📞  [Your phone number]", font_size=12, color=LIGHT_SAGE)

add_rect(s11, Inches(0), Inches(6.9), SLIDE_W, Inches(0.6), fill_color=MID_GREEN)
add_textbox(s11, Inches(0.5), Inches(6.92), Inches(12.0), Inches(0.5),
            "Marketing Plan 2026  |  TOCS — Top Owners Corporation Solution  |  portal.tocs.com.au",
            font_size=11, color=CREAM, align=PP_ALIGN.LEFT)


# ── Save ───────────────────────────────────────────────────────────────────────
output_path = "/home/user/Ocorder/marketing/TOCS-OC-Order-Portal-Marketing-Plan.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
